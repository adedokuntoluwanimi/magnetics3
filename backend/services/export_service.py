from __future__ import annotations

import base64
import io
import json
import re
import zipfile
from datetime import datetime, timezone
from typing import Any

import pandas as pd

from backend.logging_utils import log_event
from backend.models import ExportJob, ExportRequest


def utc_now() -> str:
    return datetime.now(timezone.utc).isoformat()


class ExportService:
    def __init__(self, store, storage_backend, ai_service) -> None:
        self._store = store
        self._storage = storage_backend
        self._ai_service = ai_service

    def _load_full_results(self, task: dict) -> dict:
        artifacts = (task.get("results") or {}).get("artifacts") or []
        results_artifact = next((a for a in artifacts if a.get("file_name") == "results.json"), None)
        if not results_artifact:
            raise ValueError("Results file not found. Run processing first.")
        data = self._storage.download_bytes(results_artifact["bucket"], results_artifact["object_name"])
        return json.loads(data)

    def create_export(self, task_id: str, request: ExportRequest) -> dict:
        task = self._store.get_task(task_id)
        if not task:
            raise ValueError("Task not found.")
        if not any(a.get("file_name") == "results.json" for a in ((task.get("results") or {}).get("artifacts") or [])):
            raise ValueError("Run processing before exporting.")
        full_results = self._load_full_results(task)
        project = self._store.get_project(task["project_id"])
        export_context = {
            "formats": request.formats,
            "aurora_sections": request.aurora_sections,
            "export_options": request.export_options,
        }
        aurora = self._ai_service.generate_export_report(task["project_id"], task_id, full_results, export_context)
        job = ExportJob(task_id=task_id, formats=request.formats, status="running")
        persisted = self._store.create_export_job(job.model_dump(mode="json"))
        artifacts_out = []
        for export_format in request.formats:
            artifact = self._build_artifact(project, task, full_results, export_format, aurora, request)
            if artifact:
                artifacts_out.append(artifact.model_dump(mode="json"))
        completed = self._store.update_export_job(
            persisted["id"],
            {
                "status": "completed",
                "updated_at": utc_now(),
                "details": {
                    "artifacts": artifacts_out,
                    "aurora": aurora.model_dump(mode="json"),
                    "sections": request.aurora_sections,
                    "export_options": request.export_options,
                },
            },
        )
        jobs = list(task.get("export_jobs") or [])
        jobs.append(completed)
        self._store.update_task(task_id, {"export_jobs": jobs[-10:], "updated_at": utc_now()})
        log_event("INFO", "Export completed", action="export.complete", task_id=task_id, export_job_id=completed["id"])
        return completed

    def _build_artifact(self, project: dict, task: dict, data: dict, export_format: str, aurora, request: ExportRequest):
        normalized = export_format.lower()
        builders = {
            "csv": ("application/zip", f"{self._artifact_stem(task, project)}_csv_bundle.zip", lambda: self._build_csv_bundle(project, task, data, request)),
            "geojson": ("application/zip", f"{self._artifact_stem(task, project)}_geojson_bundle.zip", lambda: self._build_geojson_bundle(project, task, data, request)),
            "kmz": ("application/zip", f"{self._artifact_stem(task, project)}_kmz_bundle.zip", lambda: self._build_kmz_bundle(project, task, data, request)),
            "kml": ("application/zip", f"{self._artifact_stem(task, project)}_kmz_bundle.zip", lambda: self._build_kmz_bundle(project, task, data, request)),
            "gdb": ("application/zip", f"{self._artifact_stem(task, project)}_gdb_bundle.zip", lambda: self._build_gdb_bundle(project, task, data, request)),
            "pdf": ("application/pdf", f"{self._artifact_stem(task, project)}.pdf", lambda: self._build_pdf(project, task, data, aurora)),
            "word": ("application/vnd.openxmlformats-officedocument.wordprocessingml.document", f"{self._artifact_stem(task, project)}.docx", lambda: self._build_docx(project, task, data, aurora)),
            "pptx": ("application/vnd.openxmlformats-officedocument.presentationml.presentation", f"{self._artifact_stem(task, project)}.pptx", lambda: self._build_pptx(project, task, data, aurora)),
            "png": ("application/zip", f"{self._artifact_stem(task, project)}_map_bundle.zip", lambda: self._build_raster_bundle(project, task, data, request)),
            "jpg": ("application/zip", f"{self._artifact_stem(task, project)}_map_bundle.zip", lambda: self._build_raster_bundle(project, task, data, request)),
            "jpeg": ("application/zip", f"{self._artifact_stem(task, project)}_map_bundle.zip", lambda: self._build_raster_bundle(project, task, data, request)),
        }
        spec = builders.get(normalized)
        if not spec:
            return None
        content_type, file_name, builder = spec
        return self._storage.upload_result(project_id=task["project_id"], task_id=task["id"], file_name=file_name, content_type=content_type, data=builder(), kind="export")

    def _artifact_stem(self, task: dict, project: dict) -> str:
        return (re.sub(r"[^A-Za-z0-9._-]+", "_", f"{task.get('name', 'task')}_{project.get('name', 'project')}").strip("._") or "gaia_export").lower()

    def _flags(self, request: ExportRequest | None) -> dict[str, bool]:
        defaults = {
            "include_corrected_field": True,
            "include_regional_field": False,
            "include_residual_field": False,
            "include_anomaly_catalogue": False,
            "include_drilling_recommendations": False,
            "include_structural_interpretation": False,
            "include_coverage_gap_analysis": False,
        }
        if request:
            defaults.update({k: bool(v) for k, v in (request.export_options or {}).items() if k in defaults})
        return defaults

    def _layer_specs(self, data: dict) -> list[dict[str, str]]:
        specs = [
            {"key": "surface", "slug": "corrected_field", "label": "Corrected Magnetic Field"},
            {"key": "corrected_field", "slug": "corrected_field", "label": "Corrected Magnetic Field"},
            {"key": "regional_field", "slug": "regional_field", "label": "Regional Magnetic Field"},
            {"key": "regional_surface", "slug": "regional_field", "label": "Regional Magnetic Field"},
            {"key": "regional_residual", "slug": "residual_field", "label": "Residual Magnetic Field"},
            {"key": "residual_surface", "slug": "residual_field", "label": "Residual Magnetic Field"},
            {"key": "rtp_surface", "slug": "rtp", "label": self._rtp_export_label(data)},
            {"key": "analytic_signal", "slug": "analytic_signal", "label": "Analytic Signal"},
            {"key": "tilt_derivative", "slug": "tilt_derivative", "label": "Tilt Derivative"},
            {"key": "total_gradient", "slug": "total_gradient", "label": "Total Gradient"},
            {"key": "first_vertical_derivative", "slug": "fvd", "label": "First Vertical Derivative"},
            {"key": "horizontal_derivative", "slug": "hd", "label": "Horizontal Derivative"},
            {"key": "uncertainty", "slug": "uncertainty", "label": "Uncertainty"},
            {"key": "filtered_surface", "slug": "filtered_surface", "label": "Filtered Surface"},
            {"key": "upward_continuation", "slug": "upward_continuation", "label": "Upward Continuation"},
            {"key": "downward_continuation", "slug": "downward_continuation", "label": "Downward Continuation"},
            {"key": "emag2_residual", "slug": "regional_residual_product", "label": "Regional Residual Product"},
        ]
        active, seen = [], set()
        for spec in specs:
            if data.get(spec["key"]) is not None and spec["slug"] not in seen:
                active.append(spec)
                seen.add(spec["slug"])
        return active

    def _selected_layers(self, data: dict, request: ExportRequest | None = None) -> list[dict[str, str]]:
        flags = self._flags(request)
        selected = []
        for spec in self._layer_specs(data):
            if spec["slug"] == "corrected_field" and not flags["include_corrected_field"]:
                continue
            if spec["slug"] == "regional_field" and not flags["include_regional_field"]:
                continue
            if spec["slug"] == "residual_field" and not flags["include_residual_field"]:
                continue
            selected.append(spec)
        return selected

    def _grid_rows(self, data: dict, key: str, value_field: str) -> list[dict]:
        rows = []
        for lat_row, lon_row, value_row in zip(data.get("grid_y") or [], data.get("grid_x") or [], data.get(key) or []):
            for latitude, longitude, value in zip(lat_row, lon_row, value_row):
                rows.append({"longitude": longitude, "latitude": latitude, value_field: value})
        return rows

    def _point_rows(self, points: list[dict]) -> list[dict]:
        return [dict(point) for point in (points or [])]

    def _point_value(self, point: dict):
        for key in ("magnetic", "corrected_magnetic", "predicted_magnetic", "value"):
            if point.get(key) is not None:
                return point.get(key)
        return None

    def _points_geojson(self, points: list[dict], *, value_field: str = "magnetic") -> dict:
        return {"type": "FeatureCollection", "features": [{"type": "Feature", "geometry": {"type": "Point", "coordinates": [row["longitude"], row["latitude"]]}, "properties": {**{k: v for k, v in row.items() if k not in {"longitude", "latitude"}}, value_field: row.get(value_field, self._point_value(row))}} for row in points if row.get("longitude") is not None and row.get("latitude") is not None]}

    def _traverse_geojson(self, points: list[dict]) -> dict:
        groups: dict[str, list[dict]] = {}
        for point in points or []:
            key = str(point.get("line_id") or point.get("traverse_label") or "line")
            groups.setdefault(key, []).append(point)
        features = []
        for line_id, group in groups.items():
            coords = [[p.get("longitude"), p.get("latitude")] for p in group if p.get("longitude") is not None and p.get("latitude") is not None]
            if len(coords) >= 2:
                features.append({"type": "Feature", "geometry": {"type": "LineString", "coordinates": coords}, "properties": {"line_id": line_id}})
        return {"type": "FeatureCollection", "features": features}

    def _build_kml(self, points: list[dict], *, name: str, value_field: str = "magnetic") -> str:
        placemarks = []
        for row in points or []:
            longitude = row.get("longitude")
            latitude = row.get("latitude")
            if longitude is None or latitude is None:
                continue
            value = row.get(value_field, self._point_value(row))
            placemarks.append("<Placemark>" f"<name>{value if value is not None else name}</name>" f"<description>{name}</description>" f"<Point><coordinates>{longitude},{latitude},0</coordinates></Point>" "</Placemark>")
        return "<?xml version='1.0' encoding='UTF-8'?>" "<kml xmlns='http://www.opengis.net/kml/2.2'><Document>" f"<name>{name}</name>" + "".join(placemarks) + "</Document></kml>"

    def _bundle_metadata(self, project: dict, task: dict, data: dict, bundle_name: str, manifest: list[str], request: ExportRequest | None = None) -> dict[str, Any]:
        qa = data.get("qa_report") or {}
        correction = data.get("correction_report") or {}
        return {
            "project_metadata": {"project_name": project.get("name"), "project_context": project.get("context"), "task_name": task.get("name"), "task_description": task.get("description")},
            "run_metadata": {"task_id": task.get("id"), "project_id": task.get("project_id"), "bundle_name": bundle_name, "exported_at": utc_now()},
            "scenario_type": task.get("scenario") or "automatic",
            "selected_processing_options": task.get("analysis_config") or {},
            "actual_algorithm_path": {
                "model_used": data.get("model_used") or (task.get("analysis_config") or {}).get("model") or "not specified",
                "regional_method": data.get("regional_method_used"),
                "fallback_events": qa.get("fallback_events") or [],
            },
            "applied_corrections": correction.get("applied_order") or (task.get("analysis_config") or {}).get("corrections") or [],
            "skipped_corrections": [event.get("requested") for event in (correction.get("fallbacks") or []) if event.get("requested")],
            "modelling_methods": [data.get("model_used") or (task.get("analysis_config") or {}).get("model") or "not specified"],
            "regional_method": data.get("regional_method_used"),
            "derived_layers_available": [spec["label"] for spec in self._layer_specs(data)],
            "diagnostics_summary": {"processing_quality_score": qa.get("quality_score"), "run_status": qa.get("status"), "warning_flags": qa.get("warnings") or [], "fallback_events": qa.get("fallback_events") or [], "correction_report": correction, "model_diagnostics": data.get("model_metadata") or {}, "stage_log": data.get("stage_log") or []},
            "timestamps": {"generated_at": utc_now()},
            "file_manifest": manifest,
            "export_selections_used": {"formats": request.formats if request else [], "aurora_sections": request.aurora_sections if request else [], "export_options": request.export_options if request else {}},
        }

    def _surface_key(self, slug: str) -> str:
        mapping = {
            "corrected_field": "surface",
            "regional_field": "regional_field",
            "residual_field": "regional_residual",
            "rtp": "rtp_surface",
            "analytic_signal": "analytic_signal",
            "tilt_derivative": "tilt_derivative",
            "total_gradient": "total_gradient",
            "fvd": "first_vertical_derivative",
            "hd": "horizontal_derivative",
            "uncertainty": "uncertainty",
            "filtered_surface": "filtered_surface",
            "upward_continuation": "upward_continuation",
            "downward_continuation": "downward_continuation",
            "regional_residual_product": "emag2_residual",
        }
        return mapping.get(slug, slug)

    def _grid_for_slug(self, data: dict, slug: str) -> list[dict]:
        key = self._surface_key(slug)
        if key == "regional_field" and data.get(key) is None:
            key = "regional_surface"
        if key == "regional_residual" and data.get(key) is None:
            key = "residual_surface"
        return self._grid_rows(data, key, "value") if data.get(key) is not None else []

    def _write_json(self, archive: zipfile.ZipFile, path: str, payload: Any) -> None:
        archive.writestr(path, json.dumps(payload, indent=2, default=str))

    def _write_csv(self, archive: zipfile.ZipFile, path: str, rows: list[dict]) -> None:
        frame = pd.DataFrame(rows)
        archive.writestr(path, frame.to_csv(index=False) if not frame.empty else "")

    def _build_csv_bundle(self, project: dict, task: dict, data: dict, request: ExportRequest | None = None) -> bytes:
        manifest: list[str] = []
        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
            for spec in self._selected_layers(data, request):
                rows = self._grid_for_slug(data, spec["slug"])
                if rows:
                    path = f"csv_bundle/{spec['slug']}.csv"
                    self._write_csv(archive, path, rows)
                    manifest.append(path)
            measured = self._point_rows(data.get("points") or [])
            predicted = self._point_rows(data.get("predicted_points") or [])
            if measured:
                self._write_csv(archive, "csv_bundle/measured_points.csv", measured)
                manifest.append("csv_bundle/measured_points.csv")
            if predicted:
                self._write_csv(archive, "csv_bundle/predicted_points.csv", predicted)
                manifest.append("csv_bundle/predicted_points.csv")
            if measured or predicted:
                self._write_csv(archive, "csv_bundle/line_profiles.csv", measured + predicted)
                manifest.append("csv_bundle/line_profiles.csv")
            if data.get("correction_report"):
                self._write_json(archive, "csv_bundle/correction_report.json", data.get("correction_report"))
                manifest.append("csv_bundle/correction_report.json")
            if data.get("model_metadata"):
                self._write_json(archive, "csv_bundle/model_diagnostics.json", data.get("model_metadata"))
                manifest.append("csv_bundle/model_diagnostics.json")
            metadata_path = "csv_bundle/metadata.json"
            self._write_json(archive, metadata_path, self._bundle_metadata(project, task, data, "csv_bundle", manifest + [metadata_path], request))
        return buffer.getvalue()

    def _build_geojson_bundle(self, project: dict, task: dict, data: dict, request: ExportRequest | None = None) -> bytes:
        manifest: list[str] = []
        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
            measured = self._point_rows(data.get("points") or [])
            if measured:
                self._write_json(archive, "geojson_bundle/corrected_points.geojson", self._points_geojson(measured))
                manifest.append("geojson_bundle/corrected_points.geojson")
                traverses = self._traverse_geojson(measured)
                if traverses.get("features"):
                    self._write_json(archive, "geojson_bundle/line_profiles.geojson", traverses)
                    manifest.append("geojson_bundle/line_profiles.geojson")
            predicted = self._point_rows(data.get("predicted_points") or [])
            if predicted:
                self._write_json(archive, "geojson_bundle/predicted_points.geojson", self._points_geojson(predicted, value_field="predicted_magnetic"))
                manifest.append("geojson_bundle/predicted_points.geojson")
            for spec in self._selected_layers(data, request):
                if spec["slug"] in {"corrected_field", "regional_field", "residual_field", "rtp", "analytic_signal", "tilt_derivative", "total_gradient", "uncertainty"}:
                    rows = self._grid_for_slug(data, spec["slug"])
                    if rows:
                        file_slug = "corrected_surface" if spec["slug"] == "corrected_field" else spec["slug"]
                        path = f"geojson_bundle/{file_slug}.geojson"
                        self._write_json(archive, path, self._points_geojson(rows, value_field="value"))
                        manifest.append(path)
            metadata_path = "geojson_bundle/metadata.json"
            self._write_json(archive, metadata_path, self._bundle_metadata(project, task, data, "geojson_bundle", manifest + [metadata_path], request))
        return buffer.getvalue()

    def _image_specs(self, data: dict, request: ExportRequest | None = None) -> list[dict[str, str]]:
        specs = [
            ("corrected_heatmap_b64", "corrected_heatmap.png"),
            ("corrected_contour_b64", "corrected_contour.png"),
            ("regional_heatmap_b64", "regional_heatmap.png"),
            ("regional_contour_b64", "regional_contour.png"),
            ("residual_heatmap_b64", "residual_heatmap.png"),
            ("residual_contour_b64", "residual_contour.png"),
            ("rtp_b64", "rtp.png"),
            ("analytic_signal_b64", "analytic.png"),
            ("tilt_derivative_b64", "tilt.png"),
            ("total_gradient_b64", "total_gradient.png"),
            ("first_vertical_derivative_b64", "fvd.png"),
            ("horizontal_derivative_b64", "hd.png"),
            ("uncertainty_b64", "uncertainty.png"),
            ("upward_continuation_b64", "upward_continuation.png"),
            ("downward_continuation_b64", "downward_continuation.png"),
            ("emag2_comparison_b64", "regional_residual_product.png"),
            ("line_profile_b64", "line_profile.png"),
            ("anomaly_catalogue_b64", "anomaly_catalogue.png"),
        ]
        flags = self._flags(request)
        out = []
        for key, name in specs:
            if name.startswith("regional_") and not flags["include_regional_field"]:
                continue
            if name.startswith("residual_") and not flags["include_residual_field"]:
                continue
            if name.startswith("corrected_") and not flags["include_corrected_field"]:
                continue
            if data.get(key):
                out.append({"key": key, "file_name": name})
        return out

    def _decode_b64(self, payload: str) -> bytes:
        body = payload.split(",", 1)[1] if isinstance(payload, str) and payload.startswith("data:") else payload
        return base64.b64decode(body)

    def _build_kmz_bundle(self, project: dict, task: dict, data: dict, request: ExportRequest | None = None) -> bytes:
        manifest: list[str] = []
        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
            measured = self._point_rows(data.get("points") or [])
            kml_map = [
                ("corrected", measured),
                ("regional", self._grid_for_slug(data, "regional_field")),
                ("residual", self._grid_for_slug(data, "residual_field")),
                ("rtp", self._grid_for_slug(data, "rtp")),
                ("analytic_signal", self._grid_for_slug(data, "analytic_signal")),
                ("tilt_derivative", self._grid_for_slug(data, "tilt_derivative")),
            ]
            for slug, rows in kml_map:
                if rows:
                    path = f"kml_bundle/{slug}.kml"
                    archive.writestr(path, self._build_kml(rows, name=slug.replace("_", " ").title(), value_field="value" if slug != "corrected" else "magnetic"))
                    manifest.append(path)
            overlay_map = {
                "corrected_overlay.png": data.get("corrected_heatmap_b64"),
                "regional_overlay.png": data.get("regional_heatmap_b64"),
                "residual_overlay.png": data.get("residual_heatmap_b64"),
                "rtp_overlay.png": data.get("rtp_b64"),
                "analytic_overlay.png": data.get("analytic_signal_b64"),
                "tilt_overlay.png": data.get("tilt_derivative_b64"),
            }
            for file_name, payload in overlay_map.items():
                if payload:
                    path = f"kml_bundle/overlays/{file_name}"
                    archive.writestr(path, self._decode_b64(payload))
                    manifest.append(path)
            metadata_path = "kml_bundle/metadata.json"
            self._write_json(archive, metadata_path, self._bundle_metadata(project, task, data, "kml_bundle", manifest + [metadata_path], request))
        return buffer.getvalue()

    def _build_gdb_bundle(self, project: dict, task: dict, data: dict, request: ExportRequest | None = None) -> bytes:
        manifest: list[str] = []
        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
            feature_map = {
                "corrected_feature_class": self._point_rows(data.get("points") or []),
                "regional_feature_class": self._grid_for_slug(data, "regional_field"),
                "residual_feature_class": self._grid_for_slug(data, "residual_field"),
                "rtp_feature_class": self._grid_for_slug(data, "rtp"),
                "analytic_feature_class": self._grid_for_slug(data, "analytic_signal"),
                "tilt_feature_class": self._grid_for_slug(data, "tilt_derivative"),
                "total_gradient_feature_class": self._grid_for_slug(data, "total_gradient"),
                "uncertainty_feature_class": self._grid_for_slug(data, "uncertainty"),
                "predicted_points_feature_class": self._point_rows(data.get("predicted_points") or []),
            }
            for name, rows in feature_map.items():
                if rows:
                    path = f"gdb_bundle/{name}.geojson"
                    self._write_json(archive, path, self._points_geojson(rows, value_field="value"))
                    manifest.append(path)
            metadata_path = "gdb_bundle/metadata.json"
            self._write_json(archive, metadata_path, self._bundle_metadata(project, task, data, "gdb_bundle", manifest + [metadata_path], request))
        return buffer.getvalue()

    def _build_raster_bundle(self, project: dict, task: dict, data: dict, request: ExportRequest | None = None) -> bytes:
        manifest: list[str] = []
        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
            for spec in self._image_specs(data, request):
                path = f"maps/{spec['file_name']}"
                archive.writestr(path, self._decode_b64(data[spec["key"]]))
                manifest.append(path)
            metadata_path = "maps/metadata.json"
            self._write_json(archive, metadata_path, self._bundle_metadata(project, task, data, "maps", manifest + [metadata_path], request))
        return buffer.getvalue()

    def _rtp_export_label(self, data: dict) -> str:
        fallback_events = (data.get("qa_report") or {}).get("fallback_events") or []
        if any("rtp" in str(event.get("requested") or "").lower() for event in fallback_events):
            return "Low-latitude fallback transform"
        return "Reduction to Pole (RTP)"

    def _brand(self) -> dict[str, Any]:
        return {
            "navy": "#0E2A47",
            "teal": "#138A72",
            "cyan": "#2F9FB3",
            "sand": "#F4F1EA",
            "panel": "#F7F9FB",
            "text": "#1A2433",
            "muted": "#5D6B7A",
            "border": "#D8E0E8",
            "warn": "#C96A28",
            "danger": "#B23A48",
        }

    def _status_label(self, data: dict) -> str:
        status = str((data.get("qa_report") or {}).get("status") or "unknown").lower()
        mapping = {
            "valid": "Valid",
            "valid_with_warnings": "Valid with Warnings",
            "degraded": "Degraded",
            "failed": "Failed",
        }
        return mapping.get(status, status.replace("_", " ").title() or "Unknown")

    def _quality_check_payload(self, payload: dict[str, Any], *, slide_mode: bool = False) -> dict[str, Any]:
        blocked_patterns = ("{", "}", "[Insert", "TODO", "TBD", "placeholder", "undefined", "null")
        cleaned_sections = []
        seen_bodies: set[str] = set()
        seen_bullets: set[str] = set()
        quality_checks = list(payload.get("quality_checks") or [])
        for section in payload.get("sections") or []:
            title = (section.get("title") or "").strip()
            body = (section.get("body") or "").strip()
            bullets = [str(item).strip() for item in (section.get("bullets") or []) if str(item).strip()]
            if slide_mode:
                bullets = bullets[:4]
                body = body[:280].strip()
            else:
                bullets = bullets[:6]
                body = body[:1200].strip()
            if not title:
                continue
            if not bullets and not body and not (section.get("visual_refs") or []) and not (section.get("table_rows") or []):
                continue
            normalized_body = re.sub(r"\s+", " ", body.lower())
            normalized_bullets = " | ".join(re.sub(r"\s+", " ", bullet.lower()) for bullet in bullets)
            if normalized_body and normalized_body in seen_bodies:
                quality_checks.append(f"Removed repeated section body in {title}.")
                body = ""
            if normalized_body:
                seen_bodies.add(normalized_body)
            if normalized_bullets and normalized_bullets in seen_bullets:
                quality_checks.append(f"Trimmed repeated layer bullets in {title}.")
                bullets = bullets[:1]
            if normalized_bullets:
                seen_bullets.add(normalized_bullets)
            if any(pattern in body for pattern in blocked_patterns if pattern not in {"{", "}"}):
                body = re.sub(r"\[Insert[^\]]+\]", "", body).strip()
                quality_checks.append(f"Removed unresolved placeholder text from {title}.")
            if "latitude,longitude" in body.lower() or "magnetic\n" in body.lower() or body.count(",") > 24:
                body = "Acquisition detail was summarised into metrics and findings so raw survey rows do not spill into the narrative body."
                quality_checks.append(f"Summarised raw row-like content in {title}.")
            if body.startswith("{") and body.endswith("}"):
                body = "Diagnostics were summarised into formatted tables and status panels during rendering."
                quality_checks.append(f"Converted raw JSON-like body content in {title}.")
            interpretation = str(section.get("interpretation") or "").strip()
            implication = str(section.get("implication") or "").strip()
            if interpretation and implication and re.sub(r"\s+", " ", interpretation.lower()) == re.sub(r"\s+", " ", implication.lower()):
                implication = ""
                quality_checks.append(f"Removed repeated interpretation/implication text in {title}.")
            if section.get("type") == "layer":
                linked_text = " ".join([body, section.get("interpretation") or "", section.get("implication") or "", " ".join(bullets)]).lower()
                if not any(term in linked_text for term in ("regional", "residual", "corrected", "derivative", "analytic", "rtp")):
                    quality_checks.append(f"Added cross-layer interpretation reminder to {title}.")
                    bullets.append("Interpret this layer alongside the corrected, regional, residual, or derivative products rather than in isolation.")
                if ("warning" not in linked_text and "fallback" not in linked_text and "confidence" not in linked_text and "uncertainty" not in linked_text):
                    quality_checks.append(f"Added QA-aware caution to {title}.")
                    section["callouts"] = list(section.get("callouts") or []) + ["Confidence should be judged against QA warnings, fallback events, and uncertainty support."]
            if section.get("type") == "conclusions":
                conclusion_text = " ".join([body, " ".join(bullets)]).lower()
                if not any(term in conclusion_text for term in ("corrected", "regional", "residual", "derivative", "qa", "uncertainty", "fallback")):
                    quality_checks.append("Conclusions were missing explicit linkage to results and QA context.")
                    bullets.append("Conclusions should follow from the corrected, regional, residual, derivative, and QA evidence rather than from isolated map descriptions.")
            if section.get("type") in {"recommendations", "drilling_recommendations"}:
                recommendation_text = " ".join(bullets + ([body] if body else [])).lower()
                generic_markers = ("verify before drilling", "additional work may be needed", "consider follow-up", "more data may help")
                if any(marker in recommendation_text for marker in generic_markers):
                    quality_checks.append(f"Generic recommendation wording detected in {title}.")
                if not any(term in recommendation_text for term in ("tie", "crossover", "inclination", "declination", "station", "support", "uncertainty", "fallback", "level", "prediction")):
                    quality_checks.append(f"Recommendations in {title} were not clearly tied to limitations.")
                    bullets.append("Tie each recommendation to a specific limitation such as leveling support, RTP inputs, support density, or uncertainty.")
            section["title"] = title
            section["body"] = body
            section["interpretation"] = interpretation
            section["implication"] = implication
            section["bullets"] = bullets
            section["callouts"] = [str(item).strip() for item in (section.get("callouts") or []) if str(item).strip()][:4]
            section["table_rows"] = [row for row in (section.get("table_rows") or []) if isinstance(row, dict)][:8]
            cleaned_sections.append(section)
        payload["sections"] = cleaned_sections
        payload["quality_checks"] = quality_checks
        return payload

    def _ordered_sections(self, payload: dict[str, Any], *, slide_mode: bool = False) -> dict[str, Any]:
        order = {
            "project_overview": 30,
            "data_summary": 40,
            "processing_workflow": 50,
            "key_findings": 60,
            "layer": 70,
            "data_quality": 80,
            "structural_interpretation": 90,
            "anomaly_catalogue": 95,
            "coverage_gap_analysis": 100,
            "drilling_recommendations": 110,
            "conclusions": 120,
            "quality_check": 200,
        }
        layer_order = {
            "Corrected Magnetic Field": 0,
            "Regional Magnetic Field": 1,
            "Residual Magnetic Field": 2,
            "Reduction to Pole (RTP)": 3,
            "Low-latitude fallback transform": 4,
            "Analytic Signal": 5,
            "Tilt Derivative": 6,
            "Total Gradient": 7,
            "First Vertical Derivative": 8,
            "Horizontal Derivative": 9,
            "Uncertainty": 10,
        }
        sections = payload.get("sections") or []
        payload["sections"] = sorted(
            sections,
            key=lambda section: (
                order.get(section.get("type") or "", 150),
                layer_order.get(section.get("title") or "", 50),
                section.get("title") or "",
            ),
        )
        if slide_mode:
            payload["sections"] = payload["sections"][:14]
        return payload

    def _visual_payload_map(self, data: dict) -> dict[str, str]:
        return {
            "corrected_heatmap.png": "corrected_heatmap_b64",
            "corrected_contour.png": "corrected_contour_b64",
            "regional_heatmap.png": "regional_heatmap_b64",
            "regional_contour.png": "regional_contour_b64",
            "residual_heatmap.png": "residual_heatmap_b64",
            "residual_contour.png": "residual_contour_b64",
            "line_profile.png": "line_profile_b64",
            "anomaly_catalogue.png": "anomaly_catalogue_b64",
        }

    def _grid_key_for_visual(self, visual_ref: str) -> str | None:
        return {
            "surface.png": "surface",
            "regional_surface.png": "regional_surface",
            "residual_surface.png": "residual_surface",
            "rtp.png": "rtp_surface",
            "analytic_signal.png": "analytic_signal",
            "tilt_derivative.png": "tilt_derivative",
            "total_gradient.png": "total_gradient",
            "first_vertical_derivative.png": "first_vertical_derivative",
            "horizontal_derivative.png": "horizontal_derivative",
            "uncertainty.png": "uncertainty",
            "upward_continuation.png": "upward_continuation",
            "downward_continuation.png": "downward_continuation",
            "regional_residual_product.png": "emag2_residual",
        }.get(visual_ref)

    def _visual_title(self, visual_ref: str) -> str:
        return visual_ref.replace(".png", "").replace("_", " ").title()

    def _figure_image_bytes(self, data: dict, visual_ref: str) -> bytes:
        payload_key = self._visual_payload_map(data).get(visual_ref)
        if payload_key and data.get(payload_key):
            return self._decode_b64(data[payload_key])
        grid_key = self._grid_key_for_visual(visual_ref)
        if grid_key and data.get(grid_key) is not None:
            return self._render_grid_preview(data.get(grid_key), self._visual_title(visual_ref))
        if visual_ref == "line_profile.png" and data.get("points"):
            return self._render_line_profile_preview(data.get("points") or [])
        return self._render_placeholder_figure(self._visual_title(visual_ref), "Figure preview was not persisted, so a styled reference card was generated.")

    def _render_grid_preview(self, grid: Any, title: str) -> bytes:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt
        import numpy as np

        arr = np.asarray(grid, dtype=float)
        fig, ax = plt.subplots(figsize=(8.4, 4.8), dpi=170)
        image = ax.imshow(arr, cmap="viridis", origin="lower", aspect="auto")
        ax.set_title(title, fontsize=12, fontweight="bold")
        ax.set_xticks([])
        ax.set_yticks([])
        cbar = fig.colorbar(image, ax=ax, fraction=0.03, pad=0.02)
        cbar.ax.tick_params(labelsize=8)
        fig.tight_layout()
        buffer = io.BytesIO()
        fig.savefig(buffer, format="png", bbox_inches="tight", facecolor="white")
        plt.close(fig)
        return buffer.getvalue()

    def _render_line_profile_preview(self, points: list[dict]) -> bytes:
        import matplotlib
        matplotlib.use("Agg")
        import matplotlib.pyplot as plt

        xs = []
        ys = []
        for idx, point in enumerate(points):
            xs.append(float(point.get("along_line_m", idx)))
            ys.append(float(self._point_value(point) or 0.0))
        fig, ax = plt.subplots(figsize=(8.4, 4.2), dpi=170)
        ax.plot(xs, ys, color="#138A72", linewidth=2.2)
        ax.scatter(xs, ys, color="#0E2A47", s=10)
        ax.set_title("Line Profile", fontsize=12, fontweight="bold")
        ax.set_xlabel("Distance")
        ax.set_ylabel("nT")
        ax.grid(True, alpha=0.22)
        fig.tight_layout()
        buffer = io.BytesIO()
        fig.savefig(buffer, format="png", bbox_inches="tight", facecolor="white")
        plt.close(fig)
        return buffer.getvalue()

    def _render_placeholder_figure(self, title: str, caption: str) -> bytes:
        from PIL import Image, ImageDraw, ImageFont

        brand = self._brand()
        image = Image.new("RGB", (1400, 760), brand["sand"])
        draw = ImageDraw.Draw(image)
        title_font = ImageFont.load_default()
        body_font = ImageFont.load_default()
        draw.rounded_rectangle((36, 36, 1364, 724), radius=28, fill="white", outline=brand["border"], width=3)
        draw.rectangle((36, 36, 1364, 132), fill=brand["navy"])
        draw.text((72, 68), title, fill="white", font=title_font)
        draw.text((72, 186), caption, fill=brand["text"], font=body_font)
        buffer = io.BytesIO()
        image.save(buffer, format="PNG")
        return buffer.getvalue()

    def _document_payload(self, aurora, key: str) -> dict[str, Any]:
        report_data = getattr(aurora, "report_data", None) or {}
        payload = report_data.get(key) if isinstance(report_data, dict) else None
        if isinstance(payload, dict):
            cleaned = self._quality_check_payload(dict(payload), slide_mode=(key == "pptx"))
            return self._ordered_sections(cleaned, slide_mode=(key == "pptx"))
        return {
            "title": f"{key.upper()} Export",
            "executive_summary": getattr(aurora, "summary", ""),
            "sections": [],
            "visuals": [],
            "recommendations": getattr(aurora, "highlights", []),
            "file_manifest": [],
            "quality_checks": [],
        }

    def _render_doc_text(self, project: dict, task: dict, payload: dict[str, Any]) -> str:
        lines = [
            payload.get("title") or f"{task.get('name', 'Task')} Export",
            payload.get("subtitle") or "",
            "",
            payload.get("executive_summary") or "",
        ]
        for section in payload.get("sections") or []:
            if not section.get("include", True):
                continue
            lines.extend(["", section.get("title") or section.get("type") or "Section"])
            for bullet in section.get("bullets") or []:
                if bullet:
                    lines.append(f"- {bullet}")
            if section.get("body"):
                lines.append(section["body"])
            if section.get("caption"):
                lines.append(f"Caption: {section['caption']}")
            if section.get("interpretation"):
                lines.append(f"Interpretation: {section['interpretation']}")
            if section.get("implication"):
                lines.append(f"Implication: {section['implication']}")
        if payload.get("recommendations"):
            lines.extend(["", "Recommendations"])
            for item in payload.get("recommendations") or []:
                if item:
                    lines.append(f"- {item}")
        return "\n".join([line for line in lines if line is not None])

    def _summary_table_rows(self, task: dict, data: dict) -> list[dict[str, Any]]:
        stats = data.get("stats") or {}
        qa = data.get("qa_report") or {}
        validation = data.get("validation_summary") or {}
        return [
            {"Metric": "Scenario", "Value": task.get("scenario") or "automatic"},
            {"Metric": "Model", "Value": data.get("model_used") or (task.get("analysis_config") or {}).get("model") or "not specified"},
            {"Metric": "Point count", "Value": stats.get("point_count", 0)},
            {"Metric": "Base-station readings", "Value": validation.get("base_station_count", 0)},
            {"Metric": "QA status", "Value": self._status_label(data)},
            {"Metric": "Quality score", "Value": qa.get("quality_score", "n/a")},
            {"Metric": "Warnings", "Value": len(qa.get("warnings") or [])},
            {"Metric": "Fallbacks", "Value": len(qa.get("fallback_events") or [])},
        ]

    def _append_quality_section(self, payload: dict[str, Any], data: dict) -> dict[str, Any]:
        if not payload.get("quality_checks"):
            return payload
        payload["sections"] = list(payload.get("sections") or []) + [
            {
                "type": "quality_check",
                "title": "Export Quality Check",
                "include": True,
                "bullets": payload.get("quality_checks"),
                "body": "The renderer removed unresolved placeholders, repeated boilerplate, and raw dump patterns before finalising the export.",
                "visual_refs": [],
                "figure_title": "",
                "caption": "",
                "interpretation": "",
                "implication": "",
                "layout": "callout",
                "callouts": [f"QA status: {self._status_label(data)}"],
                "table_rows": [],
                "speaker_notes": "",
            }
        ]
        return payload

    def _section_family(self, section: dict[str, Any]) -> str:
        mapping = {
            "project_overview": "Context",
            "data_summary": "Context",
            "processing_workflow": "Workflow",
            "key_findings": "Findings",
            "layer": "Results",
            "data_quality": "Reliability",
            "structural_interpretation": "Interpretation",
            "anomaly_catalogue": "Interpretation",
            "coverage_gap_analysis": "Reliability",
            "conclusions": "Conclusions",
            "drilling_recommendations": "Next Steps",
            "quality_check": "Appendix",
        }
        return mapping.get(section.get("type") or "", "Results")

    def _paired_visual_sections(self, payload: dict[str, Any]) -> list[dict[str, Any]]:
        sections = payload.get("sections") or []
        by_title = {section.get("title"): section for section in sections}
        pairs = []
        if by_title.get("Corrected Magnetic Field") and by_title.get("Regional Magnetic Field"):
            pairs.append({
                "title": "Corrected vs Regional Magnetic Field",
                "left": by_title["Corrected Magnetic Field"],
                "right": by_title["Regional Magnetic Field"],
                "footer": "Corrected field versus broad background trend",
            })
        if by_title.get("Residual Magnetic Field") and by_title.get("Analytic Signal"):
            pairs.append({
                "title": "Residual vs Analytic Signal",
                "left": by_title["Residual Magnetic Field"],
                "right": by_title["Analytic Signal"],
                "footer": "Residual anomaly expression versus source-edge sharpening",
            })
        return pairs

    def _build_pdf(self, project: dict, task: dict, data: dict, aurora) -> bytes:
        payload = self._append_quality_section(self._document_payload(aurora, "pdf"), data)
        try:
            from reportlab.lib import colors
            from reportlab.lib.enums import TA_CENTER
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib.units import cm
            from reportlab.platypus import Image, KeepTogether, PageBreak, Paragraph, SimpleDocTemplate, Spacer, Table, TableStyle

            brand = self._brand()
            styles = getSampleStyleSheet()
            styles.add(ParagraphStyle(name="GaiaTitle", parent=styles["Title"], fontName="Helvetica-Bold", fontSize=24, textColor=colors.HexColor(brand["navy"]), spaceAfter=18))
            styles.add(ParagraphStyle(name="GaiaSubtitle", parent=styles["Heading2"], fontName="Helvetica", fontSize=12, textColor=colors.HexColor(brand["muted"]), spaceAfter=10))
            styles.add(ParagraphStyle(name="GaiaHeading", parent=styles["Heading1"], fontName="Helvetica-Bold", fontSize=15, textColor=colors.HexColor(brand["navy"]), spaceBefore=14, spaceAfter=8))
            styles.add(ParagraphStyle(name="GaiaBody", parent=styles["BodyText"], fontName="Helvetica", fontSize=9.5, leading=14, textColor=colors.HexColor(brand["text"]), spaceAfter=8))
            styles.add(ParagraphStyle(name="GaiaCaption", parent=styles["BodyText"], fontName="Helvetica-Oblique", fontSize=8.5, textColor=colors.HexColor(brand["muted"]), spaceAfter=10))
            styles.add(ParagraphStyle(name="GaiaChip", parent=styles["BodyText"], fontName="Helvetica-Bold", fontSize=9.5, textColor=colors.white, alignment=TA_CENTER))

            story = [
                Spacer(1, 3.2 * cm),
                Paragraph(payload.get("title") or f"{task.get('name', 'Task')} Technical Export", styles["GaiaTitle"]),
                Paragraph(payload.get("subtitle") or project.get("name", ""), styles["GaiaSubtitle"]),
            ]
            chip = Table([[Paragraph(self._status_label(data), styles["GaiaChip"])]], colWidths=[5.0 * cm])
            chip.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), colors.HexColor(brand["teal"])), ("BOX", (0, 0), (-1, -1), 0, colors.HexColor(brand["teal"])), ("TOPPADDING", (0, 0), (-1, -1), 6), ("BOTTOMPADDING", (0, 0), (-1, -1), 6)]))
            story.extend([chip, Spacer(1, 0.35 * cm), Paragraph(payload.get("executive_summary") or "", styles["GaiaBody"])])
            summary_rows = [["Metric", "Value"]] + [[row["Metric"], str(row["Value"])] for row in self._summary_table_rows(task, data)]
            summary_table = Table(summary_rows, colWidths=[5.2 * cm, 8.2 * cm])
            summary_table.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(brand["navy"])), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white), ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"), ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor(brand["border"])), ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(brand["panel"])])]))
            story.extend([Spacer(1, 0.25 * cm), summary_table, PageBreak()])

            for section in payload.get("sections") or []:
                story.append(Paragraph(section.get("title") or "Section", styles["GaiaHeading"]))
                for callout in (section.get("callouts") or [])[:3]:
                    panel = Table([[callout]], colWidths=[14.5 * cm])
                    panel.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, -1), colors.HexColor(brand["sand"])), ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor(brand["border"])), ("TOPPADDING", (0, 0), (-1, -1), 6), ("BOTTOMPADDING", (0, 0), (-1, -1), 6)]))
                    story.extend([panel, Spacer(1, 0.15 * cm)])
                for bullet in (section.get("bullets") or [])[:5]:
                    story.append(Paragraph(f"• {bullet}", styles["GaiaBody"]))
                if section.get("table_rows"):
                    keys = list(section["table_rows"][0].keys())
                    rows = [keys] + [[str(row.get(key, "")) for key in keys] for row in section["table_rows"][:6]]
                    table = Table(rows, colWidths=[14.5 * cm / max(1, len(keys))] * max(1, len(keys)))
                    table.setStyle(TableStyle([("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(brand["navy"])), ("TEXTCOLOR", (0, 0), (-1, 0), colors.white), ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor(brand["border"])), ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor(brand["panel"])])]))
                    story.extend([table, Spacer(1, 0.15 * cm)])
                if section.get("body"):
                    story.append(Paragraph(section["body"], styles["GaiaBody"]))
                visual_story = []
                for visual_ref in (section.get("visual_refs") or [])[:2]:
                    image = Image(io.BytesIO(self._figure_image_bytes(data, visual_ref)), width=14.2 * cm, height=7.4 * cm)
                    visual_story.append(image)
                    caption = " ".join([part for part in [section.get("figure_title") or self._visual_title(visual_ref), section.get("caption"), section.get("interpretation"), section.get("implication")] if part])
                    if caption:
                        visual_story.append(Paragraph(caption, styles["GaiaCaption"]))
                if visual_story:
                    story.append(KeepTogether(visual_story))
                story.append(Spacer(1, 0.2 * cm))

            if payload.get("recommendations"):
                story.extend([PageBreak(), Paragraph("Recommendations and Next Steps", styles["GaiaHeading"])])
                for item in payload.get("recommendations")[:6]:
                    story.append(Paragraph(f"• {item}", styles["GaiaBody"]))

            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=A4, leftMargin=1.7 * cm, rightMargin=1.7 * cm, topMargin=1.5 * cm, bottomMargin=1.5 * cm)

            def _page(canvas, _doc):
                canvas.setFont("Helvetica", 8)
                canvas.setFillColor(colors.HexColor(brand["muted"]))
                canvas.drawString(1.7 * cm, 1.0 * cm, f"{project.get('name', '')} | {task.get('name', '')}")
                canvas.drawRightString(A4[0] - 1.7 * cm, 1.0 * cm, f"GAIA Magnetics | {self._status_label(data)}")

            doc.build(story, onFirstPage=_page, onLaterPages=_page)
            return buffer.getvalue()
        except Exception:
            text = self._render_doc_text(project, task, payload)
            safe = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
            stream = f"BT /F1 11 Tf 50 780 Td 14 TL ({safe[:5000].replace(chr(10), ') Tj T* (')}) Tj ET"
            objects = [
                "1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj",
                "2 0 obj << /Type /Pages /Kids [3 0 R] /Count 1 >> endobj",
                "3 0 obj << /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >> endobj",
                f"4 0 obj << /Length {len(stream.encode('latin-1', errors='replace'))} >> stream\n{stream}\nendstream endobj",
                "5 0 obj << /Type /Font /Subtype /Type1 /BaseFont /Helvetica >> endobj",
            ]
            buffer = io.BytesIO()
            buffer.write(b"%PDF-1.4\n")
            offsets = [0]
            for obj in objects:
                offsets.append(buffer.tell())
                buffer.write(obj.encode("latin-1", errors="replace"))
                buffer.write(b"\n")
            xref = buffer.tell()
            buffer.write(f"xref\n0 {len(offsets)}\n".encode("latin-1"))
            buffer.write(b"0000000000 65535 f \n")
            for offset in offsets[1:]:
                buffer.write(f"{offset:010d} 00000 n \n".encode("latin-1"))
            buffer.write(f"trailer << /Size {len(offsets)} /Root 1 0 R >>\nstartxref\n{xref}\n%%EOF".encode("latin-1"))
            return buffer.getvalue()

    def _build_docx(self, project: dict, task: dict, data: dict, aurora) -> bytes:
        payload = self._append_quality_section(self._document_payload(aurora, "docx"), data)
        try:
            from docx import Document
            from docx.enum.section import WD_SECTION
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.oxml import OxmlElement
            from docx.oxml.ns import qn
            from docx.shared import Cm, Pt, RGBColor

            brand = self._brand()
            document = Document()
            section = document.sections[0]
            section.top_margin = Cm(1.8)
            section.bottom_margin = Cm(1.5)
            section.left_margin = Cm(2.0)
            section.right_margin = Cm(2.0)
            section.header.paragraphs[0].text = f"{project.get('name', '')} | GAIA Magnetics"
            section.footer.paragraphs[0].text = self._status_label(data)

            styles = document.styles
            styles["Normal"].font.name = "Aptos"
            styles["Normal"].font.size = Pt(10.5)
            styles["Heading 1"].font.name = "Aptos Display"
            styles["Heading 1"].font.size = Pt(16)
            styles["Heading 1"].font.bold = True
            styles["Heading 1"].font.color.rgb = RGBColor.from_string(brand["navy"].replace("#", ""))

            title = document.add_paragraph()
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = title.add_run(payload.get("title") or f"{task.get('name', 'Task')} Report")
            run.font.name = "Aptos Display"
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = RGBColor.from_string(brand["navy"].replace("#", ""))
            if payload.get("subtitle"):
                subtitle = document.add_paragraph()
                subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
                subrun = subtitle.add_run(payload["subtitle"])
                subrun.font.size = Pt(11)
                subrun.font.color.rgb = RGBColor.from_string(brand["muted"].replace("#", ""))

            summary_table = document.add_table(rows=0, cols=2)
            summary_table.style = "Table Grid"
            for row in self._summary_table_rows(task, data):
                cells = summary_table.add_row().cells
                cells[0].text = str(row["Metric"])
                cells[1].text = str(row["Value"])
            if payload.get("executive_summary"):
                document.add_paragraph(payload["executive_summary"])
            document.add_page_break()

            toc_p = document.add_paragraph()
            fld_begin = OxmlElement("w:fldChar")
            fld_begin.set(qn("w:fldCharType"), "begin")
            instr = OxmlElement("w:instrText")
            instr.set(qn("xml:space"), "preserve")
            instr.text = 'TOC \\o "1-3" \\h \\z \\u'
            fld_end = OxmlElement("w:fldChar")
            fld_end.set(qn("w:fldCharType"), "end")
            toc_p._p.append(fld_begin)
            toc_p._p.append(instr)
            toc_p._p.append(fld_end)
            document.add_page_break()

            for section in payload.get("sections") or []:
                if not section.get("include", True):
                    continue
                document.add_heading(section.get("title") or section.get("type") or "Section", level=1)
                if section.get("callouts"):
                    callout = document.add_table(rows=len(section["callouts"][:3]), cols=1)
                    callout.style = "Table Grid"
                    for idx, item in enumerate(section["callouts"][:3]):
                        callout.cell(idx, 0).text = item
                if section.get("table_rows"):
                    keys = list(section["table_rows"][0].keys())
                    table = document.add_table(rows=1, cols=len(keys))
                    table.style = "Table Grid"
                    for idx, key in enumerate(keys):
                        table.rows[0].cells[idx].text = str(key)
                    for row in section["table_rows"][:6]:
                        cells = table.add_row().cells
                        for idx, key in enumerate(keys):
                            cells[idx].text = str(row.get(key, ""))
                if section.get("body"):
                    document.add_paragraph(section["body"])
                for bullet in section.get("bullets") or []:
                    if bullet:
                        document.add_paragraph(bullet, style="List Bullet")
                for visual_ref in (section.get("visual_refs") or [])[:2]:
                    document.add_picture(io.BytesIO(self._figure_image_bytes(data, visual_ref)), width=Cm(15.5))
                    caption_text = " ".join([part for part in [section.get("figure_title") or self._visual_title(visual_ref), section.get("caption"), section.get("interpretation"), section.get("implication")] if part])
                    caption = document.add_paragraph(caption_text)
                    caption.alignment = WD_ALIGN_PARAGRAPH.CENTER
            if payload.get("recommendations"):
                document.add_section(WD_SECTION.NEW_PAGE)
                document.add_heading("Recommendations and Next Steps", level=1)
                for item in payload.get("recommendations")[:6]:
                    document.add_paragraph(item, style="List Bullet")
            buffer = io.BytesIO()
            document.save(buffer)
            return buffer.getvalue()
        except Exception:
            return self._minimal_docx(self._render_doc_text(project, task, payload))

    def _minimal_docx(self, text: str) -> bytes:
        body = "".join(f"<w:p><w:r><w:t>{self._xml_escape(line)}</w:t></w:r></w:p>" for line in text.splitlines() if line.strip())
        document_xml = f"<?xml version=\"1.0\" encoding=\"UTF-8\" standalone=\"yes\"?><w:document xmlns:w=\"http://schemas.openxmlformats.org/wordprocessingml/2006/main\"><w:body>{body}<w:sectPr/></w:body></w:document>"
        types_xml = "<?xml version=\"1.0\" encoding=\"UTF-8\" standalone=\"yes\"?><Types xmlns=\"http://schemas.openxmlformats.org/package/2006/content-types\"><Default Extension=\"rels\" ContentType=\"application/vnd.openxmlformats-package.relationships+xml\"/><Default Extension=\"xml\" ContentType=\"application/xml\"/><Override PartName=\"/word/document.xml\" ContentType=\"application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml\"/></Types>"
        rels_xml = "<?xml version=\"1.0\" encoding=\"UTF-8\" standalone=\"yes\"?><Relationships xmlns=\"http://schemas.openxmlformats.org/package/2006/relationships\"><Relationship Id=\"rId1\" Type=\"http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument\" Target=\"word/document.xml\"/></Relationships>"
        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
            archive.writestr("[Content_Types].xml", types_xml)
            archive.writestr("_rels/.rels", rels_xml)
            archive.writestr("word/document.xml", document_xml)
        return buffer.getvalue()

    def _build_pptx(self, project: dict, task: dict, data: dict, aurora) -> bytes:
        payload = self._append_quality_section(self._document_payload(aurora, "pptx"), data)
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
            from pptx.util import Inches, Pt

            brand = self._brand()
            prs = Presentation()

            def paint_background(slide, fill_hex: str) -> None:
                fill = slide.background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor.from_string(fill_hex.replace("#", ""))

            def add_title(slide, text: str, color_hex: str = brand["navy"]):
                title_shape = slide.shapes.title
                if title_shape is None:
                    title_shape = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(11.6), Inches(0.6))
                title_shape.text = text
                p = title_shape.text_frame.paragraphs[0]
                p.font.bold = True
                p.font.size = Pt(24)
                p.font.color.rgb = RGBColor.from_string(color_hex.replace("#", ""))
                return title_shape

            def add_footer(slide, text: str) -> None:
                box = slide.shapes.add_textbox(Inches(0.55), Inches(6.95), Inches(11.6), Inches(0.25))
                p = box.text_frame.paragraphs[0]
                p.text = text
                p.font.size = Pt(9)
                p.font.color.rgb = RGBColor.from_string(brand["muted"].replace("#", ""))

            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            paint_background(title_slide, brand["sand"])
            add_title(title_slide, payload.get("title") or f"{task.get('name', 'Task')} Presentation")
            if len(title_slide.placeholders) > 1:
                subtitle = title_slide.placeholders[1]
            else:
                subtitle = title_slide.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(7.0), Inches(1.2))
            subtitle.text_frame.paragraphs[0].text = payload.get("subtitle") or payload.get("executive_summary") or ""
            subtitle.text_frame.paragraphs[0].font.size = Pt(14)
            status_shape = title_slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(5.85), Inches(2.8), Inches(0.55))
            status_shape.fill.solid()
            status_shape.fill.fore_color.rgb = RGBColor.from_string(brand["teal"].replace("#", ""))
            status_shape.line.color.rgb = RGBColor.from_string(brand["teal"].replace("#", ""))
            status_shape.text_frame.paragraphs[0].text = self._status_label(data)
            status_shape.text_frame.paragraphs[0].font.size = Pt(13)
            status_shape.text_frame.paragraphs[0].font.bold = True
            status_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            add_footer(title_slide, f"{project.get('name', '')} | GAIA Magnetics")

            summary_slide = prs.slides.add_slide(prs.slide_layouts[5])
            paint_background(summary_slide, "FFFFFF")
            add_title(summary_slide, "Executive Summary")
            left = summary_slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(6.7), Inches(4.9))
            left_frame = left.text_frame
            takeaway_items = []
            for section in payload.get("sections") or []:
                if section.get("type") == "key_findings":
                    takeaway_items.extend(section.get("bullets") or [])
            takeaway_items = takeaway_items[:3] or (payload.get("recommendations") or [])[:3]
            intro = left_frame.paragraphs[0]
            intro.text = (payload.get("executive_summary") or "")[:160]
            intro.font.size = Pt(18)
            intro.font.bold = True
            for item in takeaway_items:
                p = left_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(15)
            right = summary_slide.shapes.add_textbox(Inches(7.7), Inches(1.2), Inches(4.0), Inches(4.8))
            right_frame = right.text_frame
            for idx, row in enumerate(self._summary_table_rows(task, data)[:5]):
                p = right_frame.paragraphs[0] if idx == 0 else right_frame.add_paragraph()
                p.text = f"{row['Metric']}: {row['Value']}"
                p.font.size = Pt(14)
                p.font.bold = idx < 2
            add_footer(summary_slide, "Executive summary")

            previous_family = None
            for section in payload.get("sections") or []:
                family = self._section_family(section)
                if family != previous_family:
                    divider = prs.slides.add_slide(prs.slide_layouts[5])
                    paint_background(divider, brand["navy"])
                    add_title(divider, family, "#FFFFFF")
                    divider_sub = divider.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(8.0), Inches(1.0))
                    divider_sub.text_frame.paragraphs[0].text = section.get("title") or family
                    divider_sub.text_frame.paragraphs[0].font.size = Pt(18)
                    divider_sub.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                    add_footer(divider, f"{family} section")
                    previous_family = family
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                paint_background(slide, "FFFFFF")
                add_title(slide, section.get("title") or section.get("type") or "Section")
                if section.get("visual_refs"):
                    slide.shapes.add_picture(io.BytesIO(self._figure_image_bytes(data, section["visual_refs"][0])), Inches(0.6), Inches(1.2), width=Inches(7.1), height=Inches(4.25))
                    caption = slide.shapes.add_textbox(Inches(0.6), Inches(5.55), Inches(7.1), Inches(0.55))
                    caption.text_frame.paragraphs[0].text = section.get("caption") or section.get("figure_title") or self._visual_title(section["visual_refs"][0])
                    caption.text_frame.paragraphs[0].font.size = Pt(10)
                    notes_box = slide.shapes.add_textbox(Inches(8.0), Inches(1.2), Inches(3.7), Inches(4.8))
                    notes_frame = notes_box.text_frame
                    headline = ((section.get("bullets") or [section.get("figure_title") or section.get("title") or "Key result"])[0])[:120]
                    evidence = " ".join((section.get("bullets") or [])[1:3])[:150]
                    implication = (section.get("implication") or section.get("interpretation") or "")[:150]
                    caution = ((section.get("callouts") or [section.get("interpretation") or ""])[:1] or [""])[0][:140]
                    blocks = [
                        ("Headline", headline),
                        ("Key evidence", evidence),
                        ("Implication", implication),
                        ("Caution", caution),
                    ]
                    for idx, (label, value) in enumerate([(label, value) for label, value in blocks if value]):
                        p = notes_frame.paragraphs[0] if idx == 0 else notes_frame.add_paragraph()
                        p.text = f"{label}: {value}"
                        p.font.size = Pt(13 if idx else 14)
                        p.font.bold = idx == 0
                else:
                    body_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.25), Inches(10.8), Inches(5.2))
                    body_frame = body_box.text_frame
                    for idx, bullet in enumerate((section.get("bullets") or [])[:3]):
                        p = body_frame.paragraphs[0] if idx == 0 else body_frame.add_paragraph()
                        p.text = bullet
                        p.font.size = Pt(18 if idx == 0 else 16)
                    if section.get("body"):
                        p = body_frame.add_paragraph()
                        p.text = section["body"][:180]
                        p.font.size = Pt(14)
                notes_text = "\n".join([section.get("speaker_notes") or "", section.get("body") or "", section.get("interpretation") or "", section.get("implication") or ""]).strip()
                if notes_text:
                    slide.notes_slide.notes_text_frame.text = notes_text
                add_footer(slide, section.get("title") or "Section")

            for pair in self._paired_visual_sections(payload):
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                paint_background(slide, "FFFFFF")
                add_title(slide, pair["title"])
                left_visual = (pair["left"].get("visual_refs") or [None])[0]
                right_visual = (pair["right"].get("visual_refs") or [None])[0]
                if left_visual:
                    slide.shapes.add_picture(io.BytesIO(self._figure_image_bytes(data, left_visual)), Inches(0.55), Inches(1.25), width=Inches(5.35), height=Inches(3.35))
                if right_visual:
                    slide.shapes.add_picture(io.BytesIO(self._figure_image_bytes(data, right_visual)), Inches(6.15), Inches(1.25), width=Inches(5.35), height=Inches(3.35))
                lower = slide.shapes.add_textbox(Inches(0.7), Inches(4.85), Inches(10.8), Inches(1.25))
                tf = lower.text_frame
                p1 = tf.paragraphs[0]
                p1.text = (pair["left"].get("interpretation") or "")[:160]
                p1.font.size = Pt(14)
                p2 = tf.add_paragraph()
                p2.text = (pair["right"].get("interpretation") or "")[:160]
                p2.font.size = Pt(14)
                add_footer(slide, pair["footer"])

            warnings = (data.get("qa_report") or {}).get("warnings") or []
            fallbacks = (data.get("qa_report") or {}).get("fallback_events") or []
            if warnings or fallbacks:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                paint_background(slide, "FFFFFF")
                add_title(slide, "Limitations and Reliability")
                box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(10.8), Inches(4.8))
                tf = box.text_frame
                items = warnings + [event.get("reason", "") for event in fallbacks]
                for idx, item in enumerate([item for item in items if item][:6]):
                    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    p.text = item
                    p.font.size = Pt(16)
                add_footer(slide, "Reliability and limitations")

            if payload.get("recommendations"):
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                paint_background(slide, brand["sand"])
                add_title(slide, "Recommendations and Next Steps")
                box = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(10.7), Inches(4.7))
                tf = box.text_frame
                for idx, item in enumerate(payload.get("recommendations")[:5]):
                    p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                    p.text = item
                    p.font.size = Pt(18)
                add_footer(slide, "Recommendations")
            buffer = io.BytesIO()
            prs.save(buffer)
            return buffer.getvalue()
        except Exception:
            return self._minimal_pptx(payload)

    def _minimal_pptx(self, payload: dict[str, Any]) -> bytes:
        slides = [payload.get("title") or "Presentation"] + [section.get("title") or "Section" for section in payload.get("sections") or [] if section.get("include", True)]
        buffer = io.BytesIO()
        with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
            archive.writestr("[Content_Types].xml", "<Types xmlns='http://schemas.openxmlformats.org/package/2006/content-types'></Types>")
            archive.writestr("_rels/.rels", "<Relationships xmlns='http://schemas.openxmlformats.org/package/2006/relationships'></Relationships>")
            archive.writestr("ppt/presentation.xml", "<p:presentation xmlns:p='http://schemas.openxmlformats.org/presentationml/2006/main'></p:presentation>")
            archive.writestr("ppt/slides/notes.txt", "\n".join(slides))
        return buffer.getvalue()

    def _xml_escape(self, text: str) -> str:
        return text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
